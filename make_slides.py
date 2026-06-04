"""Generate the WorldCupGenAI presentation deck (8 slides, ~6 minutes)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand colors
NAVY = RGBColor(0x1E, 0x4B, 0x8F)
BLUE = RGBColor(0x00, 0x55, 0xD4)
LIGHT_BG = RGBColor(0xF0, 0xF5, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
GREY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xE5, 0x39, 0x35)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_filled_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, font_size=18, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    for i, line in enumerate(text.split("\n")):
        run_para = p if i == 0 else tf.add_paragraph()
        run_para.alignment = align
        run = run_para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_bullets(slide, x, y, w, h, items, font_size=18, color=DARK_TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(8)
    return box


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_link(slide, x, y, w, h, label, url, font_size=14, color=None,
             bold=False, align=PP_ALIGN.CENTER, prefix=""):
    """Text box with an optional plain prefix + a clickable hyperlink."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if prefix:
        r0 = p.add_run()
        r0.text = prefix
        r0.font.size = Pt(font_size)
        r0.font.bold = bold
        r0.font.name = "Calibri"
        if color:
            r0.font.color.rgb = color
    link = p.add_run()
    link.text = label
    link.font.size = Pt(font_size)
    link.font.bold = bold
    link.font.name = "Calibri"
    if color:
        link.font.color.rgb = color
    link.hyperlink.address = url
    return box


def slide_header(slide, title, subtitle=None):
    add_filled_rect(slide, 0, 0, prs.slide_width, Inches(0.55), NAVY)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(12), Inches(0.4),
             title, font_size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
                 subtitle, font_size=13, color=GREY_TEXT)


# ===================== SLIDE 1: TITLE =====================
s = prs.slides.add_slide(BLANK)
add_filled_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
         "WORLD CUP GENAI HACKATHON  ·  TRACK 1",
         font_size=14, bold=True, color=RGBColor(0xFF, 0xD7, 0x00),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.2),
         "WorldCupGPT",
         font_size=68, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8),
         "A LangChain agent for World Cup history, stats, and match predictions",
         font_size=22, color=RGBColor(0xE0, 0xE8, 0xF5),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
         "Neenu Bonny + team  ·  IE 5373 Applied Generative AI",
         font_size=16, color=RGBColor(0xC0, 0xCC, 0xE0),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
         "Prof. Mohammad Dehghani  ·  Northeastern University",
         font_size=14, color=RGBColor(0xC0, 0xCC, 0xE0),
         align=PP_ALIGN.CENTER)
add_link(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         label="github.com/Neenu1234/WorldCupGenAI",
         url="https://github.com/Neenu1234/WorldCupGenAI",
         font_size=12, color=RGBColor(0xA0, 0xB0, 0xD0))
add_notes(s,
    "Hello, we are presenting WorldCupGPT, our Track 1 submission for the "
    "World Cup GenAI Hackathon. It is a LangChain powered chatbot for "
    "questions about World Cup history, team statistics, and match predictions. "
    "(15 seconds)")

# ===================== SLIDE 2: PROBLEM =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "1.  The Problem", "What Track 1 asks us to build")

add_text(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.6),
         "A conversational chatbot that does THREE things:",
         font_size=22, bold=True, color=DARK_TEXT)

for i, (icon, label, desc) in enumerate([
    ("1", "History Q&A",
     "Answer natural language questions about past World Cup matches"),
    ("2", "Team statistics",
     "Look up titles, win rate, top scorers for any national team"),
    ("3", "Match predictions",
     "Generate a preview and predicted outcome for any two teams"),
]):
    y = Inches(2.2 + i * 1.4)
    add_filled_rect(s, Inches(0.6), y, Inches(0.8), Inches(0.8), BLUE)
    add_text(s, Inches(0.6), y + Inches(0.1), Inches(0.8), Inches(0.6),
             icon, font_size=32, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y, Inches(11), Inches(0.5),
             label, font_size=22, bold=True, color=NAVY)
    add_text(s, Inches(1.7), y + Inches(0.5), Inches(11), Inches(0.5),
             desc, font_size=16, color=DARK_TEXT)

add_text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
         "Our answer: a single ReAct agent with 4 tools, conversation "
         "memory, and a custom Streamlit UI.",
         font_size=16, bold=True, color=ACCENT)
add_notes(s,
    "Track 1 asks for a chatbot that handles three jobs: history Q&A, "
    "team stats, and match predictions. Most teams will use one RAG chain. "
    "We chose to go further and build a ReAct agent that can route between "
    "four specialized tools. (30 seconds)")

# ===================== SLIDE 3: ARCHITECTURE =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "2.  Architecture", "ReAct agent + 4 tools + 2 memory layers + Streamlit UI")

# Pipeline diagram as a vertical flow of boxes
boxes = [
    ("User question", BLUE),
    ("ReAct Agent (gpt-4o)  +  ConversationBufferMemory  +  UserPreferences", NAVY),
    ("4 Tools:  WorldCupRAG  |  MatchPredictor  |  TeamStats  |  MatchGoals", GREEN),
    ("Final answer (with cited evidence + limitations)", BLUE),
    ("Streamlit chat UI (live reasoning trace + 5 chart types)", NAVY),
]

for i, (text, color) in enumerate(boxes):
    y = Inches(1.4 + i * 1.05)
    add_filled_rect(s, Inches(1.5), y, Inches(10.3), Inches(0.75), color)
    add_text(s, Inches(1.5), y + Inches(0.15), Inches(10.3), Inches(0.5),
             text, font_size=16, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        arrow = s.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(6.4), y + Inches(0.78),
            Inches(0.5), Inches(0.27),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREY_TEXT
        arrow.line.fill.background()

add_text(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
         "Built end to end in one notebook + reproducible Streamlit demo",
         font_size=14, color=GREY_TEXT, align=PP_ALIGN.CENTER)
add_notes(s,
    "Here is the architecture. The user types a question, our ReAct agent "
    "decides which of four tools to call, the tool returns an answer plus "
    "evidence and limitations, and the Streamlit UI renders it with a chart "
    "and a live reasoning trace. Two memory layers run in parallel: "
    "conversation history and user preferences. (45 seconds)")

# ===================== SLIDE 4: DATA & REPRODUCIBILITY =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "3.  Data & Reproducibility",
             "Public, cited, cached, validated")

add_text(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.5),
         "Dataset:  martj42  ·  International Football Results 1872 to Present",
         font_size=18, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.4),
         "kaggle.com/datasets/martj42/international-football-results-from-1872-to-2017",
         font_size=13, color=BLUE)

# Two-column stats block
add_filled_rect(s, Inches(0.6), Inches(2.3), Inches(6.0), Inches(2.3), LIGHT_BG)
add_text(s, Inches(0.8), Inches(2.4), Inches(5.6), Inches(0.4),
         "What we loaded", font_size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(2.8), Inches(5.6), Inches(1.8), [
    "49,281 international matches (1872 to 2026)",
    "964 World Cup matches filtered out for RAG",
    "Every goal from goalscorers.csv (~44K rows)",
    "Penalty shootout winners (resolves Italy 2006, Argentina 2022)",
], font_size=13)

add_filled_rect(s, Inches(6.9), Inches(2.3), Inches(6.0), Inches(2.3), LIGHT_BG)
add_text(s, Inches(7.1), Inches(2.4), Inches(5.6), Inches(0.4),
         "How we validated", font_size=15, bold=True, color=NAVY)
add_bullets(s, Inches(7.1), Inches(2.8), Inches(5.6), Inches(1.8), [
    "Schema check (required columns must exist)",
    "Date parsing with errors='coerce'",
    "Score columns coerced to numeric, bad rows dropped",
    "Cross-reference shootouts.csv for tied finals",
], font_size=13)

add_text(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.4),
         "Caching", font_size=15, bold=True, color=NAVY)
add_bullets(s, Inches(0.8), Inches(5.3), Inches(12), Inches(1.8), [
    "Data files cached locally for full offline reproducibility (~7 MB)",
    "Chroma vector store persisted to disk, no re-embedding on each run",
    "Configuration centralized in src/config.py for clean setup",
], font_size=14)
add_notes(s,
    "We chose the martj42 Kaggle dataset, which has 136K downloads and is "
    "actively maintained. We cite the URL in our README, cache the CSVs in "
    "the repo, validate the schema on load, and cross-reference penalty "
    "shootouts so titles like Italy 2006 are correctly counted. (30 seconds)")

# ===================== SLIDE 5: THE 4 TOOLS =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "4.  The 4 Tools", "Each tool returns answer + evidence + limitations")

tools = [
    ("WorldCupRAG", BLUE,
     "RAG over 964 World Cup match documents",
     ["Chroma vector search (k = 12)",
      "OpenAI text-embedding-3-small",
      "Cites match docs as evidence"]),
    ("MatchPredictor", NAVY,
     "Head to head + recent form + LLM synthesis",
     ["Pulls 10 most recent meetings",
      "Last 5 matches form per team",
      "Predicted score with reasoning"]),
    ("TeamStats", GREEN,
     "Structured stats lookup from pandas",
     ["Titles, win rate, top scorers",
      "Handles penalty shootout finals",
      "Pulls top scorers from goalscorers.csv"]),
    ("MatchGoals", ACCENT,
     "Goal by goal breakdown for one historical match",
     ["Inputs: TeamA vs TeamB YYYY",
      "Returns scorer, minute, own goal, penalty",
      "Powers the goal timeline chart"]),
]

for i, (name, color, sub, bullets) in enumerate(tools):
    row, col = i // 2, i % 2
    x = Inches(0.6 + col * 6.3)
    y = Inches(1.3 + row * 2.85)
    add_filled_rect(s, x, y, Inches(6.0), Inches(2.65), LIGHT_BG)
    add_filled_rect(s, x, y, Inches(0.18), Inches(2.65), color)
    add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(5.6), Inches(0.5),
             name, font_size=20, bold=True, color=color)
    add_text(s, x + Inches(0.3), y + Inches(0.55), Inches(5.6), Inches(0.4),
             sub, font_size=13, color=GREY_TEXT)
    add_bullets(s, x + Inches(0.3), y + Inches(1.0), Inches(5.6), Inches(1.6),
                bullets, font_size=12)

add_notes(s,
    "These are our four tools. Each one is a Python function the agent can "
    "call. They all return a standard ToolOutput dataclass: answer plus "
    "evidence plus limitations. The agent picks one tool per question based "
    "on the tool descriptions. (45 seconds)")

# ===================== SLIDE 6: LIVE DEMO =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "5.  Live Demo", "Five clicks cover the whole project")

add_text(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
         "Demo flow (also available as a recorded screen capture)",
         font_size=15, color=GREY_TEXT)

flow = [
    ("1", "Who won the 2014 World Cup final?",
     "WorldCupRAG retrieves the 12 closest matches, cites the 2014 final."),
    ("2", "im a brazil and france fan, can i get their history",
     "Sidebar saves BOTH teams. Fallback chart shows Brazil vs France head to head."),
    ("3", "How many World Cup titles does Brazil have?",
     "TeamStats returns 5 titles (1958 to 2002), correctly counts penalty wins."),
    ("4", "Show goals from Germany vs Brazil 2014",
     "MatchGoals returns the 7-1 thrashing; goal timeline chart renders below."),
    ("5", "Surprise me  (sidebar button)",
     "Random query fires; demonstrates the polish."),
]

for i, (n, q, exp) in enumerate(flow):
    y = Inches(1.9 + i * 0.95)
    add_filled_rect(s, Inches(0.6), y, Inches(0.6), Inches(0.7), NAVY)
    add_text(s, Inches(0.6), y + Inches(0.13), Inches(0.6), Inches(0.5),
             n, font_size=22, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y, Inches(11.3), Inches(0.4),
             q, font_size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.4), y + Inches(0.4), Inches(11.3), Inches(0.4),
             exp, font_size=12, color=GREY_TEXT)
add_notes(s,
    "Now I will run a quick live demo. Five clicks cover all four tools, "
    "the memory, and the visuals. Watch the sidebar on the left and the "
    "agent reasoning panel on the right while I type. (2 to 2.5 minutes)")

# ===================== SLIDE 7: INNOVATION =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "6.  Innovation Beyond Base",
             "Where we went past the rubric requirements")

innovations = [
    ("Live agent reasoning panel",
     "Streamlit sidebar shows Thought / Action / Observation in real time. "
     "An explainability overlay on every answer."),
    ("5 chart types, auto-rendered by query intent",
     "Head to head (all time + recent), World Cup titles ranking, top "
     "scorers, team record (wins / draws / losses), goal by goal match "
     "timeline. Smart triggers pick the right chart for each question."),
    ("Multi-team UserPreferences memory",
     "Detects one OR multiple favorite teams + era focus from natural "
     "phrasings. Sidebar shows the live state and lets users speak with "
     "pronouns like 'their' across the chat."),
    ("Penalty shootout handling",
     "Most teams miss it; we cross reference shootouts.csv so Italy 2006 "
     "and Argentina 2022 are correctly attributed."),
    ("Suggested questions sidebar + Surprise me button",
     "Onboarding for new users, ready made demo prompts for judges."),
    ("Two-layer scope guardrails",
     "Agent refuses off-topic questions (weather, news, etc.) without calling "
     "any tool. UI also suppresses charts on refusals so the visual never "
     "contradicts a 'cannot answer' response."),
]

for i, (title, body) in enumerate(innovations):
    y = Inches(1.2 + i * 0.95)
    add_filled_rect(s, Inches(0.6), y, Inches(0.25), Inches(0.8), ACCENT)
    add_text(s, Inches(1.0), y, Inches(11.7), Inches(0.4),
             title, font_size=15, bold=True, color=NAVY)
    add_text(s, Inches(1.0), y + Inches(0.38), Inches(11.7), Inches(0.55),
             body, font_size=12, color=DARK_TEXT)
add_notes(s,
    "Quick run through what makes us different. The explainability panel, "
    "the auto-rendered charts, and the preferences memory are the three "
    "features judges should remember. (45 seconds)")

# ===================== SLIDE 8: REFLECTIONS =====================
s = prs.slides.add_slide(BLANK)
slide_header(s, "7.  Reflections & What's Next",
             "What worked, what we'd improve")

add_filled_rect(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(3.5), LIGHT_BG)
add_text(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.5),
         "What worked", font_size=18, bold=True, color=GREEN)
add_bullets(s, Inches(0.8), Inches(1.95), Inches(5.6), Inches(2.7), [
    "ReAct agent routed cleanly between 4 tools",
    "Tool descriptions were the highest-leverage prompt",
    "Streamlit + sidebar buttons = effortless demo flow",
    "Caching the Chroma store made re-runs instant",
    "Per-tool limitations satisfied the grounded reasoning rubric",
], font_size=13)

add_filled_rect(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(3.5), LIGHT_BG)
add_text(s, Inches(7.1), Inches(1.4), Inches(5.6), Inches(0.5),
         "What we'd add next", font_size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(7.1), Inches(1.95), Inches(5.6), Inches(2.7), [
    "StatsBomb event data → shot maps, pressure events",
    "Multi-team comparison (3+ teams side by side)",
    "Live match data API for 2026 fixtures",
    "Cross-era predictions (1970 Brazil vs 2022 Argentina)",
    "PDF export of match reports",
], font_size=13)

add_filled_rect(s, Inches(0.6), Inches(5.1), Inches(12.3), Inches(1.8), NAVY)
add_text(s, Inches(0.6), Inches(5.3), Inches(12.3), Inches(0.6),
         "Biggest lesson",
         font_size=16, bold=True, color=RGBColor(0xFF, 0xD7, 0x00),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.0),
         "An agent's behavior depends more on what we write (tool descriptions "
         "and ReAct prompt) than on the LLM itself. The model is the engine, "
         "the prompt and tools are the steering wheel.",
         font_size=15, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER)

add_link(s, Inches(0.6), Inches(7.0), Inches(12.3), Inches(0.4),
         prefix="Thank you  ·  Repo: ",
         label="github.com/Neenu1234/WorldCupGenAI",
         url="https://github.com/Neenu1234/WorldCupGenAI",
         font_size=14, bold=True, color=GREY_TEXT)
add_notes(s,
    "To close: ReAct agents are powerful with very little setup. The most "
    "important lesson we took away is that tool descriptions and the agent "
    "prompt drive behavior more than the model choice does. Thank you, "
    "questions welcome. (30 seconds)")

# ===================== SAVE =====================
out_path = "/Users/neenubonny/Downloads/WorldCupGenAI/WorldCupGenAI_Slides.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
